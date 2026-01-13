#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF格式转换功能
"""

import os
import fitz  # PyMuPDF
from PIL import Image


def pdf_to_word(input_path: str, output_path: str, progress_callback=None):
    """
    PDF转Word文档
    
    Args:
        input_path: 输入PDF路径
        output_path: 输出Word路径
        progress_callback: 进度回调函数
    """
    try:
        from pdf2docx import Converter
    except ImportError:
        raise ImportError("请安装 pdf2docx: pip install pdf2docx")
    
    if progress_callback:
        progress_callback(10)
    
    cv = Converter(input_path)
    
    if progress_callback:
        progress_callback(30)
    
    cv.convert(output_path)
    
    if progress_callback:
        progress_callback(90)
    
    cv.close()
    
    if progress_callback:
        progress_callback(100)
    
    return f"转换完成！已保存到 {output_path}"


def pdf_to_excel(input_path: str, output_path: str, progress_callback=None):
    """
    PDF转Excel表格（提取表格数据）
    
    Args:
        input_path: 输入PDF路径
        output_path: 输出Excel路径
        progress_callback: 进度回调函数
    """
    try:
        import pdfplumber
        from openpyxl import Workbook
    except ImportError:
        raise ImportError("请安装 pdfplumber 和 openpyxl")
    
    if progress_callback:
        progress_callback(10)
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    
    current_row = 1
    
    with pdfplumber.open(input_path) as pdf:
        total_pages = len(pdf.pages)
        
        for i, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            
            for table in tables:
                if table:
                    for row in table:
                        for col_idx, cell in enumerate(row):
                            ws.cell(row=current_row, column=col_idx + 1, 
                                   value=cell if cell else "")
                        current_row += 1
                    current_row += 1  # 表格间空一行
            
            # 如果没有表格，提取文本
            if not tables:
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        ws.cell(row=current_row, column=1, value=line)
                        current_row += 1
                    current_row += 1
            
            if progress_callback:
                progress_callback(10 + int((i + 1) / total_pages * 80))
    
    wb.save(output_path)
    
    if progress_callback:
        progress_callback(100)
    
    return f"转换完成！已保存到 {output_path}"


def pdf_to_ppt(input_path: str, output_path: str, progress_callback=None):
    """
    PDF转PowerPoint
    
    Args:
        input_path: 输入PDF路径
        output_path: 输出PPT路径
        progress_callback: 进度回调函数
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    
    if progress_callback:
        progress_callback(10)
    
    doc = fitz.open(input_path)
    prs = Presentation()
    
    # 设置幻灯片尺寸
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    total_pages = len(doc)
    
    for i, page in enumerate(doc):
        # 渲染页面为图片
        mat = fitz.Matrix(2.0, 2.0)  # 2x缩放提高质量
        pix = page.get_pixmap(matrix=mat)
        
        # 保存临时图片
        temp_img = f"/tmp/pdf_page_{i}.png"
        pix.save(temp_img)
        
        # 添加幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加图片
        slide.shapes.add_picture(
            temp_img, Inches(0), Inches(0), 
            width=prs.slide_width, height=prs.slide_height
        )
        
        # 删除临时文件
        os.remove(temp_img)
        
        if progress_callback:
            progress_callback(10 + int((i + 1) / total_pages * 85))
    
    doc.close()
    prs.save(output_path)
    
    if progress_callback:
        progress_callback(100)
    
    return f"转换完成！已保存到 {output_path}"


def pdf_to_images(input_path: str, output_dir: str, dpi: int = 150, 
                  format: str = "png", progress_callback=None):
    """
    PDF转图片
    
    Args:
        input_path: 输入PDF路径
        output_dir: 输出目录
        dpi: 图片分辨率
        format: 图片格式 (png, jpeg)
        progress_callback: 进度回调函数
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    doc = fitz.open(input_path)
    total_pages = len(doc)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    
    # 计算缩放矩阵
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    
    output_files = []
    
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat)
        
        output_path = os.path.join(output_dir, f"{base_name}_page{i+1}.{format}")
        
        if format == "jpeg":
            # JPEG需要先保存为PNG再转换
            temp_path = f"/tmp/temp_page_{i}.png"
            pix.save(temp_path)
            img = Image.open(temp_path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            img.save(output_path, 'JPEG', quality=95)
            os.remove(temp_path)
        else:
            pix.save(output_path)
        
        output_files.append(output_path)
        
        if progress_callback:
            progress_callback(int((i + 1) / total_pages * 100))
    
    doc.close()
    
    return f"转换完成！生成了 {len(output_files)} 张图片到 {output_dir}"


def word_to_pdf(input_path: str, output_path: str, progress_callback=None):
    """
    Word转PDF
    
    Args:
        input_path: 输入Word路径
        output_path: 输出PDF路径
        progress_callback: 进度回调函数
    """
    import platform
    import subprocess
    
    if progress_callback:
        progress_callback(10)
    
    system = platform.system()
    
    if system == "Darwin":  # macOS
        # 使用 LibreOffice 或 Pages
        try:
            subprocess.run([
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "--headless", "--convert-to", "pdf",
                "--outdir", os.path.dirname(output_path) or ".",
                input_path
            ], check=True, capture_output=True)
        except FileNotFoundError:
            # 尝试使用 soffice 命令
            try:
                subprocess.run([
                    "soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", os.path.dirname(output_path) or ".",
                    input_path
                ], check=True, capture_output=True)
            except FileNotFoundError:
                raise RuntimeError("请安装 LibreOffice 以支持 Word 转 PDF")
        
        # 重命名输出文件
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        temp_output = os.path.join(
            os.path.dirname(output_path) or ".",
            f"{base_name}.pdf"
        )
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)
            
    elif system == "Windows":
        try:
            import comtypes.client
            word = comtypes.client.CreateObject('Word.Application')
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(input_path))
            doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = PDF
            doc.Close()
            word.Quit()
        except ImportError:
            # 使用 LibreOffice
            try:
                subprocess.run([
                    "soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", os.path.dirname(output_path) or ".",
                    input_path
                ], check=True, capture_output=True)
            except FileNotFoundError:
                raise RuntimeError("请安装 Microsoft Word 或 LibreOffice")
    else:
        # Linux - 使用 LibreOffice
        subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", os.path.dirname(output_path) or ".",
            input_path
        ], check=True, capture_output=True)
    
    if progress_callback:
        progress_callback(100)
    
    return f"转换完成！已保存到 {output_path}"


def images_to_pdf(input_paths, output_path: str, progress_callback=None):
    """
    图片转PDF
    
    Args:
        input_paths: 输入图片路径列表
        output_path: 输出PDF路径
        progress_callback: 进度回调函数
    """
    if isinstance(input_paths, str):
        input_paths = [input_paths]
    
    if not input_paths:
        raise ValueError("请选择至少一张图片")
    
    doc = fitz.open()
    total_images = len(input_paths)
    
    for i, img_path in enumerate(input_paths):
        try:
            # 打开图片
            img = fitz.open(img_path)
            
            # 转换为PDF格式
            pdf_bytes = img.convert_to_pdf()
            img.close()
            
            # 打开转换后的PDF
            img_pdf = fitz.open("pdf", pdf_bytes)
            
            # 插入到输出文档
            doc.insert_pdf(img_pdf)
            img_pdf.close()
            
        except Exception as e:
            raise ValueError(f"无法处理图片 {img_path}: {str(e)}")
        
        if progress_callback:
            progress_callback(int((i + 1) / total_images * 90))
    
    # 保存
    doc.save(output_path)
    doc.close()
    
    if progress_callback:
        progress_callback(100)
    
    return f"转换完成！已将 {total_images} 张图片合并为 {output_path}"
